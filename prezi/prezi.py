"""
Prezi Creator

Creates PowerPoint presentations with one slide.
"""
from pptx import Presentation
from pptx.util import Inches


def create_presentation():
    """Creates an empty presentation"""
    return Presentation()


def add_style(pres):
    """Creates a first slide with the default style"""
    first_slide_layout = pres.slide_layouts[0]
    first_slide = pres.slides.add_slide(first_slide_layout)
    return first_slide


def add_title(slide, title):
    """Adds title to the slide"""
    if slide:
        slide.shapes.title.text = title
    return slide


def add_image(slide, image):
    """Adds an image to the slide"""
    if slide:
        text_box = slide.placeholders[1]
        text_box.text = "This is the sea"
        slide.shapes.add_picture(image, Inches(4), Inches(5), width=Inches(3))
    return slide


def save_presentation(pres, file_name):
    """Saves the presentation to file_name"""
    pres.save(file_name)
 

def main(title, image, file_name):
    """Runs the script producing a single slide
    presentation, with title and image
    """
    pres = create_presentation()
    slide = add_style(pres)
    slide = add_image(add_title(slide, title), image)
    save_presentation(pres, file_name)
 

if __name__ == "__main__":
    main("Gus is here!", "arabian_sea.jpeg", "prezi.pptx")